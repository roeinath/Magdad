import os
from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from pptx.dml.color import RGBColor
from pptx.util import Pt
import numpy as np
import json

def color_by_grade(grade):
    if grade < 80:   return "D10808"
    elif grade < 83: return "FF0000"
    elif grade < 87: return "FFC000"
    elif grade < 92: return "C2D69B"
    elif grade < 95: return "92D050"
    else:            return "00B050"

class PowerpointWriter:

    def __init__(self, team_directory_path, student_name):

        self.student_name = student_name
        self.file_name = student_name + '.pptx'
        self.rel_file_path = os.path.join(team_directory_path, self.file_name)

    def write_grades(self, all_final_grades):

        numeric_grade = lambda course: str.isnumeric(course['final grade'])
        all_final_grades = list(filter(numeric_grade, all_final_grades))
        all_final_grades.sort(key = lambda course: int(course['year']), reverse=True)

        # creating an object 
        ppt = Presentation("StudentTemplate.pptx")
        slide = ppt.slides[0]

        # change title to include student name
        for shp in slide.shapes:
            if shp.has_text_frame:
                shp.text_frame.paragraphs[0].text = f"{self.student_name} - ציונים אקדמיים"

        #gather all cells which can have grade inside of them
        course_name_cells = []
        grade_cells = []
        student_grade_cells = []
        for shp in slide.shapes:
            if isinstance(shp, GraphicFrame) and shp.table.cell(0,0).text_frame.text == 'קורס':
                table = shp.table
                for r in range(1, len(table.rows)): #first row is titles
                    course_name_cells.append(table.cell(r,0))
                    grade_cells.append(table.cell(r,1))
                    student_grade_cells.append(table.cell(r,2))

        #fill tables with courses
        for course, name_cell, grade_cell, mean_cell in zip(all_final_grades, course_name_cells, grade_cells, student_grade_cells):
            name_cell.text_frame.paragraphs[0].text = course['course name']
            name_cell.text_frame.paragraphs[0].font.size = Pt(12)

            if 'grade_א' in course:
                grade_cell.text_frame.paragraphs[0].text = f"{course['grade_א']}→{course['final grade']}"
                grade_cell.text_frame.paragraphs[0].font.size = Pt(10)
            else:
                grade_cell.text_frame.paragraphs[0].text = course['final grade']
                grade_cell.text_frame.paragraphs[0].font.size = Pt(12)

            try: student_mean = round(float(course['student mean']))
            except: student_mean = ""
            mean_cell.text_frame.paragraphs[0].text = f"{student_mean}"
            mean_cell.text_frame.paragraphs[0].font.size = Pt(12)

            color = color_by_grade(float(course['final grade']))
            grade_cell.fill.solid()
            grade_cell.fill.fore_color.rgb = RGBColor.from_string(color)            
        
        #search for top table
        table = None
        for shp in slide.shapes:
            if isinstance(shp, GraphicFrame) and shp.table.cell(0,0).text_frame.text == 'שנה':
                table = shp.table
        assert table is not None

        grades = np.array([float(course['final grade']) for course in all_final_grades])
        naz = np.array([float(course['naz']) for course in all_final_grades])
        year = np.array([int(course['year']) for course in all_final_grades])

        year_base = 2020 # HARD CODED - BAD!
        for i in range(0,3):
            y = year_base + i
            if len(naz[year==y]) == 0: continue
            year_mean = np.sum(grades[year==y]*naz[year==y])/np.sum(naz[year==y])
            table.cell(i,0).text_frame.paragraphs[0].text = str(y)
            table.cell(i,1).text_frame.paragraphs[0].text = f'{round(year_mean,1)}'
            table.cell(i,1).fill.solid()
            color = color_by_grade(year_mean)
            table.cell(i,1).fill.fore_color.rgb = RGBColor.from_string(color)

        total_mean = np.sum(grades*naz)/np.sum(naz)
        table.cell(3,0).text_frame.paragraphs[0].text = 'סה"כ'
        table.cell(3,1).text_frame.paragraphs[0].text = f'{round(total_mean,1)}'
        table.cell(3,1).fill.solid()
        color = color_by_grade(total_mean)
        table.cell(3,1).fill.fore_color.rgb = RGBColor.from_string(color)

        #search for moed b table
        table = None
        for shp in slide.shapes:
            if isinstance(shp, GraphicFrame) and shp.table.cell(0,0).text_frame.text == "מועדי ב'":
                table = shp.table
        assert table is not None

        num_moed_b = sum(['grade_א' in course for course in all_final_grades])
        table.cell(0,1).text_frame.paragraphs[0].text = f'{num_moed_b}'

        ppt.save(self.rel_file_path)